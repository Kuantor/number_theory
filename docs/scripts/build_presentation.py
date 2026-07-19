"""Build the number_theory presentation (issue #20) with python-pptx.

Dark, premium theme matching docs/title_slide.jpg (navy + gold + blue).
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

TITLE_IMG = sys.argv[1]
OUT = sys.argv[2]

# palette
BG    = RGBColor(0x0E, 0x1B, 0x2A)
PANEL = RGBColor(0x17, 0x2A, 0x3E)
CODEBG= RGBColor(0x0A, 0x14, 0x20)
GOLD  = RGBColor(0xD9, 0xB1, 0x3C)
GOLDL = RGBColor(0xEB, 0xCE, 0x7A)
BLUE  = RGBColor(0x5A, 0x9B, 0xD0)
WHITE = RGBColor(0xEC, 0xF2, 0xF8)
MUTED = RGBColor(0xA6, 0xBC, 0xCE)
DIM   = RGBColor(0x6E, 0x86, 0x9B)
GREEN = RGBColor(0x74, 0xC4, 0x93)
RED   = RGBColor(0xD98A, 0x00, 0x00) if False else RGBColor(0xD9, 0x7A, 0x6A)

MONO = "Consolas"
SANS = "Calibri"
SERIF = "Cambria"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def rect(s, l, t, w, h, fill=PANEL, line=None, lw=1.25, rad=0.055, rounded=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if rounded:
        try:
            shp.adjustments[0] = rad
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def tb(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    return tf


def run(p, text, size, color, bold=False, name=SANS, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = name
    return r


def para(tf, first=False, align=PP_ALIGN.LEFT, gap=6, before=0, ls=1.05):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(gap)
    p.space_before = Pt(before)
    try:
        p.line_spacing = ls
    except Exception:
        pass
    return p


def bullet(tf, text, size=15, color=WHITE, bcolor=GOLD, bold=False, first=False,
           gap=9, sub=False, glyph=None):
    p = para(tf, first=first, gap=gap, ls=1.04)
    g = glyph if glyph else ("–  " if sub else "•  ")
    rb = run(p, g, size, bcolor, bold=True)
    run(p, text, size, color, bold=bold)
    return p


def head(s, title, kicker=None, tcolor=WHITE):
    if kicker:
        tf = tb(s, 0.7, 0.44, 12, 0.5)
        p = para(tf, first=True, gap=0)
        run(p, kicker.upper(), 14, GOLD, bold=True, name=MONO)
    tf = tb(s, 0.7, 0.82, 12, 1.0)
    p = para(tf, first=True, gap=0)
    run(p, title, 33, tcolor, bold=True, name=SANS)


def footer(s, n):
    tf = tb(s, 0.7, 7.06, 6, 0.35)
    p = para(tf, first=True, gap=0)
    run(p, "number_theory  ·  v0.1.0", 10, DIM, name=MONO)
    tf2 = tb(s, 11.8, 7.06, 0.9, 0.35)
    p2 = para(tf2, first=True, gap=0, align=PP_ALIGN.RIGHT)
    run(p2, str(n), 10, DIM, name=MONO)


def codecard(s, l, t, w, h, lines, title=None, size=13.5):
    rect(s, l, t, w, h, fill=CODEBG, line=RGBColor(0x24, 0x3A, 0x50), lw=1.0)
    yoff = 0.12
    if title:
        tf = tb(s, l + 0.22, t + 0.12, w - 0.4, 0.4)
        p = para(tf, first=True, gap=0)
        run(p, title, 12, GOLD, bold=True, name=MONO)
        yoff = 0.6
    tf = tb(s, l + 0.24, t + yoff, w - 0.45, h - yoff - 0.1)
    for i, (txt, col) in enumerate(lines):
        p = para(tf, first=(i == 0), gap=4, ls=1.15)
        run(p, txt, size, col, name=MONO)


def stat(s, l, t, w, h, big, label, bigc=GOLD):
    rect(s, l, t, w, h, fill=PANEL, line=RGBColor(0x27, 0x3D, 0x54), lw=1.0)
    tf = tb(s, l + 0.1, t + 0.14, w - 0.2, h - 0.2, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True, gap=2)
    run(p, big, 30, bigc, bold=True, name=SANS)
    p2 = para(tf, gap=0)
    run(p2, label, 12.5, MUTED, name=SANS)


# ============================================================ SLIDE 1 — TITLE
s = slide()
# near-16:9 image; fit width, center vertically on the dark bg
img_h = SW / (2752 / 1536)
s.shapes.add_picture(TITLE_IMG, Inches(0), Inches((SH - img_h) / 2),
                     width=Inches(SW), height=Inches(img_h))

# ==================================================== SLIDE 2 — OVERVIEW
s = slide()
head(s, "A tiny library with a big twist", kicker="Overview")
tf = tb(s, 0.7, 2.0, 7.2, 4.6)
bullet(tf, "A lightweight, dependency-free number-theory toolkit for Python.", first=True, gap=14)
bullet(tf, "The twist: the math ships as a compiled binary — the source code stays out of the package.", gap=14)
bullet(tf, "Clean Python API plus a  number-theory  command-line tool.", gap=14)
bullet(tf, "MVP is live (v0.1.0): GCD & LCM. Lots more math on the roadmap.", gap=14)
stat(s, 8.25, 2.0, 4.35, 1.35, "v0.1.0", "first public release")
stat(s, 8.25, 3.55, 4.35, 1.35, "0", "runtime dependencies")
stat(s, 8.25, 5.10, 4.35, 1.35, "5", "functions, arbitrary precision")
footer(s, 2)

# ==================================================== SLIDE 3 — BIG IDEA
s = slide()
head(s, "Ship the math, not the source", kicker="The core idea")
tf = tb(s, 0.7, 2.0, 6.4, 4.6)
bullet(tf, "Normally  pip install  hands people your readable .py files.", first=True, gap=13)
bullet(tf, "number_theory compiles the algorithms into a native binary instead.", gap=13)
bullet(tf, "Users get working functions; the implementation stays compiled.", gap=13)
bullet(tf, "Ideal for protecting IP — or shipping an internal tool without the source.", gap=13)
# pipeline on right
px, pw = 7.5, 5.1
steps = [(".pyx", "Cython source you write", GOLD),
         ("compile", "Cython → C → native module", BLUE),
         (".pyd", "binary that ships in the wheel", GREEN)]
yy = 2.0
for i, (a, b, c) in enumerate(steps):
    rect(s, px, yy, pw, 1.15, fill=PANEL, line=c, lw=1.4)
    tf = tb(s, px + 0.25, yy + 0.13, pw - 0.5, 0.95, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True, gap=1)
    run(p, a, 19, c, bold=True, name=MONO)
    p2 = para(tf, gap=0)
    run(p2, b, 13, MUTED)
    if i < 2:
        ar = tb(s, px + pw / 2 - 0.3, yy + 1.14, 0.6, 0.34)
        p = para(ar, first=True, gap=0, align=PP_ALIGN.CENTER)
        run(p, "▼", 15, DIM)
    yy += 1.5
footer(s, 3)

# ==================================================== SLIDE 4 — NOW vs NEXT
s = slide()
head(s, "What's shipping — and what's coming", kicker="Status")
# left card: now
rect(s, 0.7, 1.95, 5.75, 4.75, fill=PANEL, line=GREEN, lw=1.5)
tf = tb(s, 1.0, 2.15, 5.2, 0.5)
p = para(tf, first=True, gap=0)
run(p, "✓  Shipping now — v0.1.0", 18, GREEN, bold=True)
tf = tb(s, 1.0, 2.75, 5.25, 3.8)
for i, (fn, d) in enumerate([
        ("gcd(a, b)", "greatest common divisor"),
        ("lcm(a, b)", "least common multiple"),
        ("extended_gcd(a, b)", "returns g, x, y with a·x+b·y=g"),
        ("gcd_many(*n)", "GCD of many integers"),
        ("lcm_many(*n)", "LCM of many integers")]):
    p = para(tf, first=(i == 0), gap=11)
    run(p, fn + "   ", 15, WHITE, bold=True, name=MONO)
    run(p, d, 12.5, MUTED)
# right card: roadmap
rect(s, 6.85, 1.95, 5.75, 4.75, fill=PANEL, line=GOLD, lw=1.5)
tf = tb(s, 7.15, 2.15, 5.2, 0.5)
p = para(tf, first=True, gap=0)
run(p, "◷  On the roadmap", 18, GOLD, bold=True)
tf = tb(s, 7.15, 2.75, 5.3, 3.8)
for i, (fn, d) in enumerate([
        ("is_prime", "primality testing  (#3)"),
        ("primes_up_to", "Sieve of Eratosthenes  (#4)"),
        ("factorize", "integer factorization  (#5)"),
        ("euler_totient  φ", "count coprime integers  (#6)"),
        ("mod_pow / mod_inverse", "modular arithmetic  (#7)")]):
    p = para(tf, first=(i == 0), gap=11)
    run(p, fn + "   ", 15, GOLDL, bold=True, name=MONO)
    run(p, d, 12.5, MUTED)
footer(s, 4)

# ==================================================== SLIDE 5 — GCD/LCM MATH
s = slide()
head(s, "The MVP: GCD & LCM, the Euclidean way", kicker="The math · shipping")
tf = tb(s, 0.7, 2.0, 6.5, 4.6)
bullet(tf, "Euclid's trick: replace (a, b) with (b, a mod b), repeat until b = 0.", first=True, gap=12)
bullet(tf, "≈ 2,300 years old — one of the oldest algorithms still in daily use.", gap=12)
bullet(tf, "Arbitrary precision: correct for numbers of ANY size.", gap=12)
bullet(tf, "extended_gcd also returns x, y with a·x + b·y = gcd — the seed of RSA key math.", gap=12)
codecard(s, 7.5, 2.0, 5.1, 3.5, title="python", lines=[
    ("import number_theory as nt", WHITE),
    ("", WHITE),
    ("nt.gcd(12, 18)            # 6", WHITE),
    ("nt.lcm(4, 6)              # 12", WHITE),
    ("nt.gcd_many(12, 18, 24)   # 6", WHITE),
    ("nt.extended_gcd(240, 46)", WHITE),
    ("#  (2, -9, 47)", GOLDL),
])
footer(s, 5)

# ==================================================== SLIDE 6 — ROADMAP MATH GRID
s = slide()
head(s, "The math still to come", kicker="The math · roadmap")
cards = [
    ("is_prime", "Is this number prime? A fast Miller–Rabin test.", GOLD),
    ("prime sieve", "All primes up to N, via the Sieve of Eratosthenes.", GOLD),
    ("factorize", "Break a number into primes:  60 = 2·2·3·5.", GOLD),
    ("Euler's totient  φ", "How many integers below n share no factor with it.", GOLD),
    ("modular arithmetic", "mod_pow, mod_inverse — the machinery behind RSA.", GOLD),
    ("fun fact", "There are infinitely many primes — Euclid proved it ~300 BC.", BLUE),
]
gw, gh, gx, gy, gap = 3.86, 1.9, 0.7, 2.0, 0.28
for i, (t, d, c) in enumerate(cards):
    r, col = divmod(i, 3)
    x = gx + col * (gw + gap)
    y = gy + r * (gh + gap)
    rect(s, x, y, gw, gh, fill=PANEL, line=RGBColor(0x2A, 0x40, 0x57), lw=1.0)
    tf = tb(s, x + 0.22, y + 0.18, gw - 0.44, gh - 0.3)
    p = para(tf, first=True, gap=6)
    run(p, t, 16, (BLUE if c == BLUE else GOLDL), bold=True, name=(MONO if c != BLUE else SANS))
    p2 = para(tf, gap=0, ls=1.08)
    run(p2, d, 13, MUTED)
footer(s, 6)

# ==================================================== SLIDE 7 — WHEEL
s = slide()
head(s, "What is a Python wheel?", kicker="Technical detour  1 / 2")
tf = tb(s, 0.7, 2.0, 6.5, 4.6)
bullet(tf, "A wheel (.whl) is Python's built package format — really just a ZIP with a special name.", first=True, gap=13)
bullet(tf, "\"Built\" means ready to install:  pip  unzips it into place — no compiling on the user's machine.", gap=13)
bullet(tf, "The filename encodes what it fits: interpreter, version, OS, and CPU.", gap=13)
bullet(tf, "Perfect for us: the compiled .pyd rides inside the wheel, so users just  pip install.", gap=13)
# decoded filename card
rect(s, 7.5, 2.0, 5.1, 3.7, fill=PANEL, line=RGBColor(0x27, 0x3D, 0x54), lw=1.0)
tf = tb(s, 7.72, 2.16, 4.7, 0.5)
p = para(tf, first=True, gap=0)
run(p, "decoding the filename", 12, GOLD, bold=True, name=MONO)
tf = tb(s, 7.72, 2.66, 4.75, 3.0)
parts = [("number_theory", "package name"), ("0.1.0", "version"),
         ("cp314", "CPython 3.14"), ("win_amd64", "Windows, 64-bit")]
for i, (a, b) in enumerate(parts):
    p = para(tf, first=(i == 0), gap=12)
    run(p, a + "   ", 15, GOLDL, bold=True, name=MONO)
    run(p, "→  " + b, 13, MUTED)
footer(s, 7)

# ==================================================== SLIDE 8 — CYTHON
s = slide()
head(s, "Cython — and why compiling is faster", kicker="Technical detour  2 / 2")
tf = tb(s, 0.7, 2.0, 7.1, 4.6)
bullet(tf, "Cython compiles Python-like code into C, then into a native module.", first=True, gap=12)
bullet(tf, "Plain Python is interpreted step-by-step; compiled C runs straight on the CPU.", gap=12)
bullet(tf, "Add static C types and hot loops can run 10–100× faster — no per-operation interpreter tax.", gap=12)
bullet(tf, "Here we use Cython mainly to hide the source; we kept full Python int math for arbitrary precision, so today's speed-ups are modest — the fast path is open for later.", gap=12)
stat(s, 8.15, 2.15, 2.15, 1.6, "1×", "interpreted\nPython", bigc=BLUE)
stat(s, 10.45, 2.15, 2.15, 1.6, "≤100×", "typed\nCython", bigc=GOLD)
tf = tb(s, 8.15, 4.0, 4.45, 2.4)
p = para(tf, first=True, gap=0, ls=1.15)
run(p, "Rule of thumb:", 13.5, GOLD, bold=True)
p = para(tf, gap=0, ls=1.2)
run(p, "  the more your code looks like tight numeric C (typed loops, no Python objects), the bigger the win.", 13.5, MUTED)
footer(s, 8)

# ==================================================== SLIDE 9 — SOURCE-LESS
s = slide()
head(s, "How we ship without the source", kicker="Under the hood")
tf = tb(s, 0.7, 2.0, 6.4, 4.6)
bullet(tf, "_core.pyx (Cython) holds every algorithm.", first=True, gap=13)
bullet(tf, "The build compiles it to  _core.<abi>.pyd  — machine code.", gap=13)
bullet(tf, "The wheel bundles only the .pyd plus thin wrappers.", gap=13)
bullet(tf, "The .pyx and generated .c are kept out. Verified: unzip the wheel and there's no algorithm source.", gap=13)
# wheel contents card
rect(s, 7.5, 2.0, 5.1, 4.2, fill=CODEBG, line=RGBColor(0x24, 0x3A, 0x50), lw=1.0)
tf = tb(s, 7.72, 2.16, 4.7, 0.5)
p = para(tf, first=True, gap=0)
run(p, "inside the .whl", 12, GOLD, bold=True, name=MONO)
tf = tb(s, 7.72, 2.7, 4.75, 3.4)
inside = [("✓  _core.cp314-win_amd64.pyd", GREEN),
          ("✓  __init__.py   (re-exports)", WHITE),
          ("✓  __main__.py   (the CLI)", WHITE),
          ("✓  _core.pyi     (type stubs)", WHITE),
          ("✗  _core.pyx     (Cython source)", RED),
          ("✗  _core.c       (generated C)", RED)]
for i, (t, c) in enumerate(inside):
    p = para(tf, first=(i == 0), gap=10, ls=1.1)
    run(p, t, 13.5, c, name=MONO)
footer(s, 9)

# ==================================================== SLIDE 10 — INSTALL
s = slide()
head(s, "Getting it: installation", kicker="User guide")
tf = tb(s, 0.7, 2.0, 6.3, 4.6)
bullet(tf, "The release wheel targets Windows x64 + CPython 3.14.", first=True, gap=13)
bullet(tf, "Install straight from the GitHub Release.", gap=13)
bullet(tf, "Other OS or Python version? Build from source (needs a C compiler).", gap=13)
bullet(tf, "Tip: install into a virtual environment to keep things tidy.", gap=13)
codecard(s, 7.4, 2.0, 5.2, 3.9, title="powershell", size=12.5, lines=[
    ("# from the GitHub Release", DIM),
    ("pip install number_theory-0.1.0-", WHITE),
    ("  cp314-cp314-win_amd64.whl", WHITE),
    ("", WHITE),
    ("# or build it yourself", DIM),
    ("pip install build", WHITE),
    ("python -m build --wheel", WHITE),
])
footer(s, 10)

# ==================================================== SLIDE 11 — CLI & API
s = slide()
head(s, "Using it: command line & Python", kicker="User guide")
codecard(s, 0.7, 2.0, 5.85, 4.3, title="command line", size=13, lines=[
    ("number-theory gcd 12 18 24", WHITE),
    ("6", GREEN),
    ("", WHITE),
    ("number-theory lcm 4 6 8", WHITE),
    ("24", GREEN),
    ("", WHITE),
    ("number-theory egcd 240 46", WHITE),
    ("2 -9 47", GREEN),
])
codecard(s, 6.75, 2.0, 5.85, 4.3, title="python", size=13, lines=[
    ("import number_theory as nt", WHITE),
    ("", WHITE),
    ("nt.gcd(12, 18)             # 6", WHITE),
    ("nt.lcm(4, 6)               # 12", WHITE),
    ("nt.gcd_many(12, 18, 24)    # 6", WHITE),
    ("nt.extended_gcd(240, 46)", WHITE),
    ("#  (2, -9, 47)", GOLDL),
])
footer(s, 11)

# ==================================================== SLIDE 12 — TROUBLESHOOTING
s = slide()
head(s, "Gotchas we hit — so you won't", kicker="User guide")
rows = [
    ("'number-theory' is not recognized",
     "Python's Scripts folder isn't on PATH. Run  python -m number_theory,  or add it to PATH."),
    ("not a supported wheel on this platform",
     "The release wheel is Windows x64 + Python 3.14 only. On anything else, build from source."),
    ("ModuleNotFoundError: number_theory._core",
     "You're running from inside the repo — the source folder shadows the installed binary. Run from elsewhere."),
]
yy = 2.05
for i, (t, d) in enumerate(rows):
    rect(s, 0.7, yy, 11.9, 1.44, fill=PANEL, line=RGBColor(0x2A, 0x40, 0x57), lw=1.0)
    tf = tb(s, 0.95, yy + 0.16, 11.4, 1.15, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True, gap=4)
    run(p, t, 15.5, GOLDL, bold=True, name=MONO)
    p2 = para(tf, gap=0, ls=1.06)
    run(p2, d, 13.5, MUTED)
    yy += 1.62
footer(s, 12)

# ==================================================== SLIDE 13 — RELEASE NOTES
s = slide()
head(s, "Release notes — v0.1.0 “GCD & LCM MVP”", kicker="Release")
tf = tb(s, 0.7, 2.0, 7.0, 4.6)
bullet(tf, "First release: GCD & LCM via Euclid, as a compiled, source-less binary.", first=True, gap=12)
bullet(tf, "Install: pip install the Windows x64 / py3.14 wheel.", gap=12)
bullet(tf, "Inside: gcd, lcm, extended_gcd, gcd_many, lcm_many — arbitrary precision.", gap=12)
bullet(tf, "Only the compiled _core.<abi>.pyd ships — no .pyx / .c.", gap=12)
bullet(tf, "Implements issues #1, #2, #10, #11.", gap=12)
rect(s, 8.0, 2.0, 4.6, 2.9, fill=PANEL, line=RGBColor(0x27, 0x3D, 0x54), lw=1.0)
tf = tb(s, 8.22, 2.16, 4.2, 0.5)
p = para(tf, first=True, gap=0)
run(p, "release assets", 12, GOLD, bold=True, name=MONO)
tf = tb(s, 8.22, 2.72, 4.25, 2.1)
for i, (t, c) in enumerate([("⬇  number_theory-0.1.0-", WHITE),
                            ("      cp314-win_amd64.whl", WHITE),
                            ("⬇  USER_GUIDE.pdf", WHITE)]):
    p = para(tf, first=(i == 0), gap=10, ls=1.1)
    run(p, t, 13.5, c, name=MONO)
footer(s, 13)

# ==================================================== SLIDE 14 — CI
s = slide()
head(s, "Built green: CI & automation", kicker="Quality")
tf = tb(s, 0.7, 2.0, 6.2, 4.6)
bullet(tf, "GitHub Actions runs on every push and pull request.", first=True, gap=13)
bullet(tf, "Each job compiles the wheel and tests the actual compiled module.", gap=13)
bullet(tf, "A release workflow builds & attaches wheels for every platform automatically.", gap=13)
bullet(tf, "Matrix below: 3 operating systems × 4 Python versions — all green.", gap=13)
# 3x4 green matrix
oss = ["Ubuntu", "Windows", "macOS"]
pys = ["3.11", "3.12", "3.13", "3.14"]
gx, gy, cw, ch, cgap = 6.95, 2.5, 1.06, 0.82, 0.12
tf = tb(s, gx, gy - 0.5, 5.5, 0.4)
p = para(tf, first=True, gap=0)
run(p, "CI matrix — 12 / 12 passing", 13, GREEN, bold=True, name=MONO)
for r, osname in enumerate(oss):
    lab = tb(s, gx - 0.0, gy + r * (ch + cgap) + 0.16, 0.0, 0.4)
    for c, py in enumerate(pys):
        x = gx + 1.15 + c * (cw + cgap)
        y = gy + r * (ch + cgap)
        rect(s, x, y, cw, ch, fill=RGBColor(0x14, 0x2C, 0x24), line=GREEN, lw=1.1)
        t2 = tb(s, x, y + 0.08, cw, ch - 0.1, anchor=MSO_ANCHOR.MIDDLE)
        p = para(t2, first=True, gap=0, align=PP_ALIGN.CENTER)
        run(p, "✓ " + py, 12.5, GREEN, bold=True, name=MONO)
    t3 = tb(s, gx - 0.05, gy + r * (ch + cgap), 1.15, ch, anchor=MSO_ANCHOR.MIDDLE)
    p = para(t3, first=True, gap=0)
    run(p, osname, 13, MUTED, name=SANS)
footer(s, 14)

# ==================================================== SLIDE 15 — ROADMAP
s = slide()
head(s, "What's next", kicker="Roadmap")
tf = tb(s, 0.7, 2.0, 11.9, 4.6)
bullet(tf, "More math: primality, the prime sieve, factorization, Euler's φ, and modular arithmetic  (#3–#7).", first=True, gap=14)
bullet(tf, "Portable wheels for every OS via cibuildwheel — and maybe a PyPI publish so  pip install number_theory  just works.", gap=14)
bullet(tf, "Tighter binary: optionally strip docstrings and compile the wrappers too.", gap=14)
bullet(tf, "Everything is tracked as issues on the project's roadmap board.", gap=14)
footer(s, 15)

# ==================================================== SLIDE 16 — CLOSING
s = slide()
tf = tb(s, 0.7, 2.5, 12, 1.4)
p = para(tf, first=True, gap=0)
run(p, "number_theory", 60, WHITE, bold=True, name=MONO)
tf = tb(s, 0.72, 3.95, 12, 0.8)
p = para(tf, first=True, gap=0)
run(p, "gcd · lcm shipping now — the rest of the primes are on their way.", 20, GOLDL, name=SANS, italic=True)
tf = tb(s, 0.72, 5.0, 12, 0.6)
p = para(tf, first=True, gap=0)
run(p, "github.com/Kuantor/number_theory", 18, MUTED, name=MONO)
tf = tb(s, 0.72, 5.7, 12, 0.6)
p = para(tf, first=True, gap=0)
run(p, "Thanks for watching!", 16, DIM, name=SANS)

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
